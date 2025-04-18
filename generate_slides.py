from pptx import Presentation
from pptx.util import Inches, Pt

def create_comparatif_swissborg_tr_v3():
    prs = Presentation()
    blank = prs.slide_layouts[6]

    # SLIDE 1
    slide1 = prs.slides.add_slide(blank)
    tb1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    p = tb1.text_frame.add_paragraph()
    p.text = "🔍 SwissBorg vs Trade Republic : quel choix pour vos investissements ?"
    p.font.bold = True; p.font.size = Pt(30)
    tb1b = slide1.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(2))
    p2 = tb1b.text_frame.add_paragraph()
    p2.text = (
        "Frais ultra‑faibles, rendement passif élevé ou simplicité ? "
        "On compare ces deux plateformes pour vous aider à décider."
    )
    p2.font.size = Pt(20)

    # SLIDE 2
    slide2 = prs.slides.add_slide(blank)
    tb2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    p = tb2.text_frame.add_paragraph()
    p.text = "🏦 Trade Republic : Aperçu"
    p.font.bold = True; p.font.size = Pt(28)
    body2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    tf2 = body2.text_frame; tf2.word_wrap = True
    p = tf2.add_paragraph(); p.text="✅ Avantages"; p.font.bold=True; p.font.size=Pt(20)
    for line in [
        "Aucun frais de transaction ni de retrait d’ATM au‑delà de 100 €.",
        "Frais de 1 € par ordre, déjà inclus dans le prix.",
        "2,5 % d’intérêt annuel sur liquidités jusqu’à 50 000 €, versés chaque mois.",
        "Investissement dès 1 € en actions, ETF, obligations ou crypto.",
        "Ouverture de compte 100 % en ligne en moins de 10 minutes."
    ]:
        p = tf2.add_paragraph(); p.text=f"→ {line}"; p.level=1; p.font.size=Pt(18)
    p = tf2.add_paragraph(); p.text="⚠️ Points à considérer"; p.font.bold=True; p.font.size=Pt(20)
    for line in [
        "Pas de rendement automatique sur liquidités au‑delà de 50 000 €.",
        "Modèle « payment for order flow » parfois critiqué.",
        "Service client uniquement en ligne, sans guichet physique."
    ]:
        p = tf2.add_paragraph(); p.text=f"→ {line}"; p.level=1; p.font.size=Pt(18)

    # SLIDE 3
    slide3 = prs.slides.add_slide(blank)
    tb3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    p = tb3.text_frame.add_paragraph()
    p.text = "💰 SwissBorg : Aperçu"
    p.font.bold = True; p.font.size = Pt(28)
    body3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    tf3 = body3.text_frame; tf3.word_wrap = True
    p = tf3.add_paragraph(); p.text="✅ Avantages"; p.font.bold=True; p.font.size=Pt(20)
    for line in [
        "Smart Yield sur plus de 20 cryptos, APY jusqu’à 37,6 % sur USDC.",
        "Retraits possibles sous 24 h sans aucun blocage.",
        "Trois profils de risque pour ajuster votre rendement.",
        "Démarrage dès 10 € sans minimum de dépôt.",
        "Communauté active et support en direct."
    ]:
        p = tf3.add_paragraph(); p.text=f"→ {line}"; p.level=1; p.font.size=Pt(18)
    p = tf3.add_paragraph(); p.text="⚠️ Points à considérer"; p.font.bold=True; p.font.size=Pt(20)
    for line in [
        "Rendements variables selon conditions de marché.",
        "Nécessite une compréhension de base des cryptos.",
        "Support parfois moins réactif en période de forte affluence."
    ]:
        p = tf3.add_paragraph(); p.text=f"→ {line}"; p.level=1; p.font.size=Pt(18)
    p = tf3.add_paragraph(); p.text="👉 Bonus : Jusqu’à 50 $ offerts avec mon lien !" ; p.level=1; p.font.size=Pt(18); p.font.bold=True

    # SLIDE 4
    slide4 = prs.slides.add_slide(blank)
    tb4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    p = tb4.text_frame.add_paragraph(); p.text="🚀 Conclusion / Avis final"; p.font.bold=True; p.font.size=Pt(28)
    body4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    tf4 = body4.text_frame; tf4.word_wrap = True
    for line in [
        "Trade Republic est parfait si vous cherchez la simplicité et des coûts quasi nuls pour actions et ETF.",
        "SwissBorg se démarque si vous souhaitez des revenus passifs élevés sur crypto tout en gardant votre liberté.",
        "Chacune a ses forces : à vous de définir vos priorités."
    ]:
        p = tf4.add_paragraph(); p.text=f"→ {line}"; p.level=0; p.font.size=Pt(18)
    p = tf4.add_paragraph(); p.text="👉 Testez SwissBorg et profitez de votre bonus de 50 $ !"; p.level=0; p.font.size=Pt(18); p.font.bold=True

    prs.save("Comparatif_SwissBorg_vs_TradeRepublic_v3.pptx")

if __name__ == "__main__":
    create_comparatif_swissborg_tr_v3()
